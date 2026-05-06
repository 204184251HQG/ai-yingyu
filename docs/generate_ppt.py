"""
生成演示视频用 9 张幻灯片
运行: python docs/generate_ppt.py
依赖: pip install python-pptx
输出: docs/词途_演示PPT.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── 配色（与软件一致）──
C_PRIMARY = RGBColor(0x0d, 0x94, 0x88)       # 深绿松石
C_PRIMARY_DARK = RGBColor(0x0f, 0x76, 0x6e)
C_PRIMARY_LIGHT = RGBColor(0xcc, 0xfb, 0xf1)
C_BG = RGBColor(0xf0, 0xfd, 0xfa)
C_TEXT = RGBColor(0x0f, 0x17, 0x2a)
C_SUB = RGBColor(0x47, 0x55, 0x69)
C_MUTED = RGBColor(0x94, 0xa3, 0xb8)
C_BORDER = RGBColor(0xe2, 0xe8, 0xf0)
C_BLUE = RGBColor(0x25, 0x63, 0xeb)
C_GREEN = RGBColor(0x16, 0xa3, 0x4a)
C_ORANGE = RGBColor(0xea, 0x58, 0x0c)
C_PURPLE = RGBColor(0x7c, 0x3a, 0xed)
C_PINK = RGBColor(0xdb, 0x27, 0x77)
C_RED = RGBColor(0xdc, 0x26, 0x26)
C_YELLOW = RGBColor(0xea, 0xb3, 0x08)
C_WHITE = RGBColor(0xff, 0xff, 0xff)

FONT_TITLE = "黑体"
FONT_BODY = "微软雅黑"


def add_text(slide, x, y, w, h, text, font_size=24, color=C_TEXT,
             bold=False, align=PP_ALIGN.LEFT, font=FONT_BODY,
             anchor=MSO_ANCHOR.TOP):
    """快捷添加文本框"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return tb


def add_rect(slide, x, y, w, h, fill, line=None, corner=False):
    """添加矩形（可圆角）"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                                  Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_card(slide, x, y, w, h, title, content, title_color=C_PRIMARY_DARK,
             card_bg=C_WHITE, border=C_BORDER):
    """添加带标题的卡片"""
    add_rect(slide, x, y, w, h, card_bg, line=border, corner=True)
    add_text(slide, x + 0.15, y + 0.15, w - 0.3, 0.5, title,
             font_size=18, color=title_color, bold=True, font=FONT_TITLE)
    add_text(slide, x + 0.15, y + 0.7, w - 0.3, h - 0.85, content,
             font_size=14, color=C_SUB)


def add_bg(slide, color=C_BG):
    add_rect(slide, 0, 0, 13.333, 7.5, color)


# ── 创建幻灯片 ──
prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 宽屏 1920x1080
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ════════════════════════════════════════════
# 张 1 — 封面
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.333, 7.5, C_PRIMARY)
# 装饰圆形
for cx, cy, r in [(11.5, 1.0, 1.5), (1.0, 6.5, 1.2), (12.5, 6.0, 0.8)]:
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(cx - r/2), Inches(cy - r/2),
                              Inches(r), Inches(r))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C_PRIMARY_DARK
    shp.line.fill.background()

# 中心圆形 logo
logo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.917), Inches(0.8),
                           Inches(1.5), Inches(1.5))
logo.fill.solid()
logo.fill.fore_color.rgb = C_WHITE
logo.line.fill.background()
add_text(s, 5.917, 0.8, 1.5, 1.5, "词途",
         font_size=48, color=C_PRIMARY, bold=True,
         align=PP_ALIGN.CENTER, font=FONT_TITLE,
         anchor=MSO_ANCHOR.MIDDLE)

add_text(s, 1, 2.6, 11.333, 1.2, "词途 · AI 英语单词学习闯关系统",
         font_size=48, color=C_WHITE, bold=True,
         align=PP_ALIGN.CENTER, font=FONT_TITLE)
add_text(s, 1, 3.9, 11.333, 0.6, "A word a day keeps the worry away",
         font_size=22, color=C_PRIMARY_LIGHT,
         align=PP_ALIGN.CENTER)
add_text(s, 1, 4.7, 11.333, 0.5, "—— 基于 5 个本地机器学习模型的百词斩式英语学习工具 ——",
         font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER)

# 字母装饰条
letters = [("A", C_BLUE), ("B", C_GREEN), ("C", C_ORANGE),
           ("D", C_PURPLE), ("E", C_PINK), ("F", C_PRIMARY_DARK),
           ("G", C_RED), ("H", C_YELLOW)]
total_w = len(letters) * 0.55
start_x = (13.333 - total_w) / 2
for i, (ch, clr) in enumerate(letters):
    x = start_x + i * 0.55
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(5.5),
                             Inches(0.5), Inches(0.5))
    sh.fill.solid()
    sh.fill.fore_color.rgb = clr
    sh.line.fill.background()
    add_text(s, x, 5.5, 0.5, 0.5, ch, font_size=18, color=C_WHITE,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, 1, 6.6, 11.333, 0.4,
         "作者: ＿＿＿  |  单位: ＿＿＿  |  覆盖小学/初中/高中  |  v1.0",
         font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# 张 2 — 痛点
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
# 标题条
add_rect(s, 0, 0, 13.333, 1.0, C_PRIMARY)
add_text(s, 0.5, 0.2, 12.333, 0.6, "教学痛点  |  Why we built this",
         font_size=28, color=C_WHITE, bold=True, font=FONT_TITLE)

# 三栏卡片
add_card(s, 0.5, 1.5, 4.0, 5.5, "1. 抄写式背诵效率低",
         "传统罚抄式作业机械重复\n3 天后遗忘率 > 70%\n\n艾宾浩斯曲线已多次验证：\n单纯重复书写无法突破\n短期记忆瓶颈\n\n孩子对「背单词」产生抵触\n甚至厌学英语",
         title_color=C_RED)

add_card(s, 4.667, 1.5, 4.0, 5.5, "2. App 课堂落地难",
         "百词斩 / 墨墨背单词等\n体验优秀但需:\n  · 联网\n  · 登录账号\n  · 绑定手机号\n\n校园机房无法批量部署\n广告与付费模块多\n\n低收入家庭难以负担",
         title_color=C_ORANGE)

add_card(s, 8.833, 1.5, 4.0, 5.5, "3. 教师缺乏诊断工具",
         "只能凭经验判断薄弱主题\n无法量化:\n  · 哪个孩子\n  · 哪个主题\n  · 正确率多少\n\n备课与个辅缺乏依据\n个性化教学无从谈起",
         title_color=C_PURPLE)


# ════════════════════════════════════════════
# 张 3 — 解决思路
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_rect(s, 0, 0, 13.333, 1.0, C_PRIMARY)
add_text(s, 0.5, 0.2, 12.333, 0.6, "解决思路  |  How it works",
         font_size=28, color=C_WHITE, bold=True, font=FONT_TITLE)

# 左 — 百词斩 4 题型
add_text(s, 0.5, 1.3, 6, 0.5, "① 百词斩式 4 种题型",
         font_size=22, color=C_PRIMARY_DARK, bold=True, font=FONT_TITLE)
quiz_modes = [
    ("看词选义", "apple → 苹果", C_BLUE),
    ("看义选词", "苹果 → apple", C_GREEN),
    ("听音辨词", "/'æpl/ → apple", C_ORANGE),
    ("拼写练习", "苹果 → a___", C_PURPLE),
]
for i, (title, demo, clr) in enumerate(quiz_modes):
    y = 1.95 + i * 1.15
    add_rect(s, 0.5, y, 6, 1.0, C_WHITE, line=clr, corner=True)
    add_rect(s, 0.5, y, 0.2, 1.0, clr)
    add_text(s, 0.85, y + 0.1, 5.5, 0.4, title,
             font_size=16, color=clr, bold=True, font=FONT_TITLE)
    add_text(s, 0.85, y + 0.5, 5.5, 0.4, demo,
             font_size=14, color=C_SUB)

# 右 — 5 个 ML 模型
add_text(s, 6.833, 1.3, 6, 0.5, "② 5 个本地 AI 模型",
         font_size=22, color=C_PRIMARY_DARK, bold=True, font=FONT_TITLE)
models = [
    ("RandomForest", "单词主题分类  54.4%"),
    ("SVM-RBF", "句型识别  97.1%"),
    ("Polynomial", "年级→词汇量  R²=0.992"),
    ("Gradient Boost", "单词难度  99.8%"),
    ("Ensemble GB", "综合主题  54.5%"),
]
for i, (algo, desc) in enumerate(models):
    y = 1.95 + i * 0.92
    add_rect(s, 6.833, y, 6, 0.78, C_WHITE, line=C_PRIMARY, corner=True)
    add_text(s, 7.05, y + 0.1, 2.4, 0.6, algo,
             font_size=14, color=C_PRIMARY, bold=True,
             font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 9.5, y + 0.1, 3.3, 0.6, desc,
             font_size=14, color=C_TEXT,
             anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════
# 张 4 — 技术全景
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_rect(s, 0, 0, 13.333, 1.0, C_PRIMARY)
add_text(s, 0.5, 0.2, 12.333, 0.6, "技术全景  |  Architecture",
         font_size=28, color=C_WHITE, bold=True, font=FONT_TITLE)

# 四层架构图
layers = [
    ("界面层", "CustomTkinter · 9 个页面 · 清新蓝绿配色 · 单文件 exe", C_BLUE),
    ("应用层", "题目派生器 · 学习记录 · 词库搜索+分页 · 自定义词库", C_PRIMARY),
    ("AI 层", "5 个 sklearn 模型 · 预训练 pkl 随包发行 · 启动即用", C_PURPLE),
    ("数据层", "4997 词全量词库（小714/初1496/高2787）· JSON 持久化", C_ORANGE),
]
for i, (name, desc, clr) in enumerate(layers):
    y = 1.5 + i * 1.0
    add_rect(s, 1.5, y, 10.333, 0.85, C_WHITE, line=clr, corner=True)
    add_rect(s, 1.5, y, 0.3, 0.85, clr, corner=False)
    add_text(s, 1.9, y + 0.1, 1.8, 0.65, name,
             font_size=18, color=clr, bold=True,
             font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 3.8, y + 0.1, 8, 0.65, desc,
             font_size=14, color=C_TEXT,
             anchor=MSO_ANCHOR.MIDDLE)
    if i < len(layers) - 1:
        # 箭头
        arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(6.5), Inches(y + 0.85),
                                     Inches(0.3), Inches(0.15))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = C_MUTED
        arrow.line.fill.background()

# 底部生成式 AI 注脚
add_rect(s, 0.5, 5.8, 12.333, 1.4, C_PRIMARY_LIGHT, line=C_PRIMARY, corner=True)
add_text(s, 0.7, 5.95, 12, 0.5, "★ 生成式 AI 全流程介入",
         font_size=18, color=C_PRIMARY_DARK, bold=True, font=FONT_TITLE)
add_text(s, 0.7, 6.45, 12, 0.7,
         "DeepSeek / 通义千问 / Claude  全程辅助：\n"
         "词库初稿生成  ·  特征工程头脑风暴  ·  代码骨架  ·  Bug 定位  ·  单元测试用例",
         font_size=13, color=C_PRIMARY_DARK)


# ════════════════════════════════════════════
# 张 5 — 目录
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_rect(s, 0, 0, 13.333, 1.0, C_PRIMARY)
add_text(s, 0.5, 0.2, 12.333, 0.6, "今日演示  |  Agenda",
         font_size=28, color=C_WHITE, bold=True, font=FONT_TITLE)

agenda = [
    ("01", "案例概述", "0:00 — 1:50", "痛点 · 思路 · 技术全景", C_BLUE),
    ("02", "实现功能", "1:50 — 6:50", "5 段录屏 · 闯关 · AI助手 · 句型 · 工具箱 · 词库 · 统计", C_PRIMARY),
    ("03", "应用情况", "6:50 — 7:50", "课堂场景 · 试用数据 · 影响力", C_PURPLE),
]
for i, (no, title, time, desc, clr) in enumerate(agenda):
    y = 1.7 + i * 1.7
    add_rect(s, 0.5, y, 12.333, 1.5, C_WHITE, line=clr, corner=True)
    # 编号圆形
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.25),
                             Inches(1), Inches(1))
    sh.fill.solid()
    sh.fill.fore_color.rgb = clr
    sh.line.fill.background()
    add_text(s, 0.8, y + 0.25, 1, 1, no,
             font_size=28, color=C_WHITE, bold=True, font=FONT_TITLE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, 2.1, y + 0.2, 6, 0.55, title,
             font_size=22, color=clr, bold=True, font=FONT_TITLE)
    add_text(s, 2.1, y + 0.8, 9, 0.5, desc,
             font_size=14, color=C_SUB)
    add_text(s, 9.8, y + 0.45, 3, 0.6, time,
             font_size=18, color=C_MUTED, bold=True,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════
# 张 6 — 课堂场景照片（占位）
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_rect(s, 0, 0, 13.333, 1.0, C_PRIMARY)
add_text(s, 0.5, 0.2, 12.333, 0.6, "应用场景  |  In the classroom",
         font_size=28, color=C_WHITE, bold=True, font=FONT_TITLE)

scenes = [
    ("课前 5 分钟热身", "全班大屏集体作答\n实时正确率反馈\n薄弱单词重点讲解", C_BLUE),
    ("机房自主练习", "1 人 1 机闯关\n错题进入下次抽题池\n已掌握单词不再优先", C_PRIMARY),
    ("教师备课诊断", "查看班级薄弱主题\n按主题正确率排序\n下周备课主攻薄弱点", C_ORANGE),
]
for i, (title, desc, clr) in enumerate(scenes):
    x = 0.5 + i * 4.333
    # 图片占位框
    add_rect(s, x, 1.5, 4.0, 2.5, C_BORDER, line=clr, corner=True)
    add_text(s, x, 1.5, 4.0, 2.5, "[ 课堂照片占位 ]\n（学生正面请打码）",
             font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # 标题
    add_text(s, x, 4.15, 4.0, 0.6, title,
             font_size=20, color=clr, bold=True, font=FONT_TITLE,
             align=PP_ALIGN.CENTER)
    # 描述
    add_rect(s, x, 4.85, 4.0, 2.0, C_WHITE, line=C_BORDER, corner=True)
    add_text(s, x + 0.15, 4.95, 3.7, 1.8, desc,
             font_size=13, color=C_SUB, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════
# 张 7 — 效果对比
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_rect(s, 0, 0, 13.333, 1.0, C_PRIMARY)
add_text(s, 0.5, 0.2, 12.333, 0.6, "试用效果  |  Outcome",
         font_size=28, color=C_WHITE, bold=True, font=FONT_TITLE)

# 左：模拟柱状图（占位）
chart_x, chart_y, chart_w, chart_h = 0.6, 1.5, 6.5, 5.5
add_rect(s, chart_x, chart_y, chart_w, chart_h, C_WHITE, line=C_BORDER, corner=True)
add_text(s, chart_x + 0.2, chart_y + 0.2, chart_w - 0.4, 0.5,
         "4 周对比  ·  单词测验平均分",
         font_size=16, color=C_TEXT, bold=True, font=FONT_TITLE)
# Y 轴
add_text(s, chart_x + 0.1, chart_y + 0.85, 0.4, 0.4, "100",
         font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)
add_text(s, chart_x + 0.1, chart_y + 4.6, 0.4, 0.4, "0",
         font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)
# 模拟柱状图（4 组对比）
weeks = [(1, 65, 70), (2, 68, 78), (3, 70, 85), (4, 72, 90)]
bar_w = 0.4
group_w = 1.3
chart_inner_x = chart_x + 0.7
chart_top = chart_y + 1.0
chart_bottom = chart_y + 4.7
chart_inner_h = chart_bottom - chart_top
for i, (wk, ctrl, exp) in enumerate(weeks):
    gx = chart_inner_x + i * group_w
    h_ctrl = chart_inner_h * ctrl / 100
    h_exp = chart_inner_h * exp / 100
    add_rect(s, gx, chart_bottom - h_ctrl, bar_w, h_ctrl, C_MUTED, corner=False)
    add_rect(s, gx + bar_w + 0.1, chart_bottom - h_exp, bar_w, h_exp, C_PRIMARY, corner=False)
    add_text(s, gx, chart_bottom + 0.05, bar_w * 2 + 0.1, 0.3,
             f"第{wk}周", font_size=10, color=C_SUB, align=PP_ALIGN.CENTER)
    add_text(s, gx, chart_bottom - h_ctrl - 0.4, bar_w, 0.3,
             str(ctrl), font_size=10, color=C_SUB, align=PP_ALIGN.CENTER)
    add_text(s, gx + bar_w + 0.1, chart_bottom - h_exp - 0.4, bar_w, 0.3,
             str(exp), font_size=10, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
# 图例
leg_y = chart_y + 5.0
add_rect(s, chart_x + 0.5, leg_y, 0.3, 0.2, C_MUTED)
add_text(s, chart_x + 0.9, leg_y - 0.05, 2, 0.3, "对照组（抄写法）",
         font_size=11, color=C_SUB)
add_rect(s, chart_x + 3.0, leg_y, 0.3, 0.2, C_PRIMARY)
add_text(s, chart_x + 3.4, leg_y - 0.05, 3, 0.3, "实验组（本系统）",
         font_size=11, color=C_PRIMARY, bold=True)

# 右：三个量化指标
metrics = [
    ("测验分提升", "+18", "分", C_GREEN),
    ("7天遗忘率下降", "-32", "pp", C_BLUE),
    ("学习兴趣提升", "+45", "%", C_ORANGE),
]
for i, (label, value, unit, clr) in enumerate(metrics):
    y = 1.5 + i * 1.85
    add_rect(s, 7.4, y, 5.4, 1.65, C_WHITE, line=clr, corner=True)
    add_rect(s, 7.4, y, 0.2, 1.65, clr)
    add_text(s, 7.7, y + 0.15, 5, 0.45, label,
             font_size=16, color=C_TEXT, bold=True, font=FONT_TITLE)
    add_text(s, 7.7, y + 0.6, 3, 0.9, value,
             font_size=44, color=clr, bold=True, font=FONT_TITLE)
    add_text(s, 10.5, y + 0.85, 1, 0.5, unit,
             font_size=18, color=clr, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, 0.5, 7.1, 12.333, 0.3,
         "* 数据为占位示例，请教师按实际课堂试用结果替换",
         font_size=10, color=C_MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# 张 8 — 影响力
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_rect(s, 0, 0, 13.333, 1.0, C_PRIMARY)
add_text(s, 0.5, 0.2, 12.333, 0.6, "影响力与共享  |  Outreach",
         font_size=28, color=C_WHITE, bold=True, font=FONT_TITLE)

cards = [
    ("代码开源", "GitHub 仓库公开\n所有源码 + 11 条 AI 提示词\nMarkdown 详细文档",
     "https://github.com/__/ai-yingyu", C_BLUE),
    ("国家平台共享", "已勾选共享至\n国家智慧教育公共\n服务平台 smartedu.cn",
     "smartedu.cn", C_PRIMARY),
    ("一线试用", "＿ 所学校\n＿ 个班级\n＿ 名学生在用",
     "持续扩展中...", C_ORANGE),
]
for i, (title, desc, footer, clr) in enumerate(cards):
    x = 0.5 + i * 4.333
    add_rect(s, x, 1.5, 4.0, 5.0, C_WHITE, line=clr, corner=True)
    add_rect(s, x, 1.5, 4.0, 0.8, clr, corner=True)
    add_text(s, x, 1.55, 4.0, 0.7, title,
             font_size=22, color=C_WHITE, bold=True, font=FONT_TITLE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.3, 2.6, 3.4, 2.5, desc,
             font_size=15, color=C_TEXT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # 二维码占位
    add_rect(s, x + 1.25, 5.2, 1.5, 1.0, C_BORDER, line=clr, corner=True)
    add_text(s, x + 1.25, 5.2, 1.5, 1.0, "[ 占位 ]",
             font_size=11, color=C_MUTED, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, 6.25, 4.0, 0.3, footer,
             font_size=11, color=clr, align=PP_ALIGN.CENTER, bold=True)


# ════════════════════════════════════════════
# 张 9 — 片尾
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.333, 7.5, C_PRIMARY_DARK)
# 装饰
for cx, cy, r in [(2.0, 1.5, 1.0), (11.5, 6.0, 1.5), (12.0, 1.0, 0.6)]:
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(cx - r/2), Inches(cy - r/2),
                             Inches(r), Inches(r))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C_PRIMARY
    sh.line.fill.background()

add_text(s, 1, 2.5, 11.333, 1.2, "Thank You !",
         font_size=72, color=C_WHITE, bold=True,
         align=PP_ALIGN.CENTER, font=FONT_TITLE)
add_text(s, 1, 3.8, 11.333, 0.5, "感谢观看  期待与各位英语教师交流",
         font_size=22, color=C_PRIMARY_LIGHT,
         align=PP_ALIGN.CENTER)

# 联系方式卡
add_rect(s, 3.5, 4.7, 6.333, 1.4, C_WHITE, corner=True)
add_text(s, 3.5, 4.85, 6.333, 0.4, "联系作者",
         font_size=16, color=C_PRIMARY, bold=True,
         align=PP_ALIGN.CENTER, font=FONT_TITLE)
add_text(s, 3.5, 5.3, 6.333, 0.4, "邮箱: ＿＿＿@＿＿＿  |  电话: ＿＿＿",
         font_size=14, color=C_TEXT, align=PP_ALIGN.CENTER)
add_text(s, 3.5, 5.7, 6.333, 0.4, "GitHub: github.com/__/ai-yingyu",
         font_size=14, color=C_SUB, align=PP_ALIGN.CENTER)

# AI 标注
add_text(s, 1, 6.6, 11.333, 0.4,
         "★ 本视频解说为本人录制，未使用 AI 语音合成",
         font_size=11, color=C_PRIMARY_LIGHT, align=PP_ALIGN.CENTER)
add_text(s, 1, 6.95, 11.333, 0.4,
         "★ 演示中 AI 模型识别画面已标注「AI 生成」",
         font_size=11, color=C_PRIMARY_LIGHT, align=PP_ALIGN.CENTER)


# ── 保存 ──
out_path = os.path.join(os.path.dirname(__file__), "词途_演示PPT.pptx")
prs.save(out_path)
print(f"PPT_GENERATED: {out_path}")
print(f"Slides: {len(prs.slides)}")
